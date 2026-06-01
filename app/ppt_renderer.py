import io
import math
import re
from dataclasses import dataclass
from functools import lru_cache
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.image_utils import image_size, load_image_from_source
from app.layout_engine import classify_panel, sort_panels_for_dashboard
from app.models import FigureAsset, Panel, PosterTask


@dataclass(frozen=True)
class Palette:
    name: str
    navy: RGBColor
    primary: RGBColor
    secondary: RGBColor
    accent: RGBColor
    success: RGBColor
    warning: RGBColor
    purple: RGBColor
    bg: RGBColor
    panel_bg: RGBColor
    soft: RGBColor
    border: RGBColor
    text: RGBColor
    muted: RGBColor
    white: RGBColor = RGBColor(255, 255, 255)


PALETTES: Dict[str, Palette] = {
    "academic_blue": Palette(
        "academic_blue",
        RGBColor(0, 43, 117),
        RGBColor(0, 82, 174),
        RGBColor(34, 134, 239),
        RGBColor(210, 35, 42),
        RGBColor(29, 138, 73),
        RGBColor(236, 139, 36),
        RGBColor(126, 63, 181),
        RGBColor(242, 247, 253),
        RGBColor(255, 255, 255),
        RGBColor(232, 241, 252),
        RGBColor(177, 199, 226),
        RGBColor(20, 32, 54),
        RGBColor(84, 98, 121),
    ),
    "engineering_green": Palette(
        "engineering_green",
        RGBColor(5, 73, 62),
        RGBColor(10, 125, 101),
        RGBColor(42, 157, 143),
        RGBColor(226, 92, 39),
        RGBColor(30, 142, 62),
        RGBColor(236, 157, 45),
        RGBColor(106, 76, 147),
        RGBColor(242, 249, 246),
        RGBColor(255, 255, 255),
        RGBColor(231, 244, 239),
        RGBColor(177, 211, 198),
        RGBColor(20, 43, 38),
        RGBColor(80, 101, 94),
    ),
    "warm_orange": Palette(
        "warm_orange",
        RGBColor(92, 50, 24),
        RGBColor(178, 84, 28),
        RGBColor(231, 126, 54),
        RGBColor(197, 36, 36),
        RGBColor(35, 138, 74),
        RGBColor(217, 154, 43),
        RGBColor(122, 78, 163),
        RGBColor(250, 246, 238),
        RGBColor(255, 255, 255),
        RGBColor(250, 236, 214),
        RGBColor(218, 189, 153),
        RGBColor(61, 42, 31),
        RGBColor(112, 91, 72),
    ),
    "minimal_gray": Palette(
        "minimal_gray",
        RGBColor(35, 39, 47),
        RGBColor(74, 87, 105),
        RGBColor(91, 141, 214),
        RGBColor(202, 55, 64),
        RGBColor(43, 150, 92),
        RGBColor(223, 145, 48),
        RGBColor(118, 92, 172),
        RGBColor(246, 248, 250),
        RGBColor(255, 255, 255),
        RGBColor(236, 240, 244),
        RGBColor(194, 203, 214),
        RGBColor(28, 33, 41),
        RGBColor(91, 101, 114),
    ),
}


KEYWORD_RE = re.compile(r"(\d+(?:\.\d+)?%|\d+(?:\.\d+)?\s?[xX×倍]|\d+(?:\.\d+)?|SOTA|state-of-the-art|GPT-\d(?:\.\d)?|Qwen|Gemini|Claude)", re.I)
METRIC_RE = re.compile(r"\d+(?:\.\d+)?%|\d+(?:\.\d+)?\s?[xX×倍]|\d+(?:,\d{3})+|\d+(?:\.\d+)?")


class Grid:
    def __init__(self, prs, x, y, w, h, cols=12, rows=8, gap=Inches(0.10)):
        self.x = x
        self.y = y
        self.w = w
        self.h = h
        self.cols = cols
        self.rows = rows
        self.gap = gap
        self.col_w = (w - gap * (cols - 1)) / cols
        self.row_h = (h - gap * (rows - 1)) / rows

    def box(self, c0: int, r0: int, cspan: int, rspan: int) -> Dict[str, int]:
        return {
            "x": self.x + c0 * (self.col_w + self.gap),
            "y": self.y + r0 * (self.row_h + self.gap),
            "w": self.col_w * cspan + self.gap * (cspan - 1),
            "h": self.row_h * rspan + self.gap * (rspan - 1),
        }


def clean_text(text: str, max_len: int = 120) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    if len(text) > max_len:
        return text[: max_len - 3] + "..."
    return text


def _emu_to_inches(value) -> float:
    return float(value) / 914400.0


def estimate_fit_font_size(
    text: str,
    w,
    h,
    *,
    max_size: float,
    min_size: float,
    line_height: float = 1.14,
    bold: bool = False,
) -> float:
    """Estimate a font size that fits wrapped text inside a PPTX box.

    LibreOffice's handling of PowerPoint auto-fit is inconsistent in headless
    rendering, so the renderer sizes text before writing the PPTX. The estimate
    is deliberately conservative; it prefers a slightly smaller font to clipped
    titles or over-eager ellipses.
    """

    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return max_size
    width_pt = max(_emu_to_inches(w) * 72.0, 1.0)
    height_pt = max(_emu_to_inches(h) * 72.0, 1.0)
    longest_word = max((len(part) for part in re.split(r"\s+", text)), default=1)
    avg_width = 0.55 if bold else 0.50
    for size in [max_size - 0.25 * i for i in range(int((max_size - min_size) / 0.25) + 1)]:
        chars_per_line = max(int(width_pt / max(size * avg_width, 1.0)), 1)
        line_count = max(
            math.ceil(len(text) / chars_per_line),
            math.ceil(longest_word / max(chars_per_line, 1)),
        )
        if line_count * size * line_height <= height_pt:
            return round(size, 2)
    return min_size


def adaptive_max_len(w, h, font_size: float, *, lines: Optional[int] = None) -> int:
    width_pt = max(_emu_to_inches(w) * 72.0, 1.0)
    height_pt = max(_emu_to_inches(h) * 72.0, 1.0)
    max_lines = lines or max(1, int(height_pt / max(font_size * 1.12, 1.0)))
    chars_per_line = max(int(width_pt / max(font_size * 0.50, 1.0)), 8)
    return max(24, chars_per_line * max_lines)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, x, y, w, h, fill, line=None, radius=True, line_width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    return add_shape(slide, shape_type, x, y, w, h, fill, line, line_width)


def add_textbox(slide, x, y, w, h, text="", font_size=12, color=None, bold=False, align=PP_ALIGN.LEFT, fit=False, min_font_size=6.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if fit:
        font_size = estimate_fit_font_size(text, w, h, max_size=font_size, min_size=min_font_size, bold=bold)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return box


def add_rich_textbox(slide, x, y, w, h, text, palette: Palette, font_size=9.2, bold=False, max_len=105, min_font_size=6.2):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if max_len <= 0:
        max_len = adaptive_max_len(w, h, font_size)
    text = clean_text(text, max_len)
    font_size = estimate_fit_font_size(text, w, h, max_size=font_size, min_size=min_font_size, bold=bold)
    pos = 0
    for match in KEYWORD_RE.finditer(text):
        if match.start() > pos:
            run = p.add_run()
            run.text = text[pos : match.start()]
            run.font.size = Pt(font_size)
            run.font.color.rgb = palette.text
            run.font.bold = bold
        run = p.add_run()
        run.text = match.group(0)
        run.font.size = Pt(font_size)
        run.font.color.rgb = palette.accent
        run.font.bold = True
        pos = match.end()
    if pos < len(text):
        run = p.add_run()
        run.text = text[pos:]
        run.font.size = Pt(font_size)
        run.font.color.rgb = palette.text
        run.font.bold = bold
    return box


def panel_icon(section: str) -> str:
    kind = classify_panel(section)
    return {
        "motivation": "?",
        "method": "1",
        "benchmark": "#",
        "innovation": "+",
        "results": "%",
        "takeaway": ">",
    }.get(kind, "-")


def panel_accent(section: str, palette: Palette) -> RGBColor:
    kind = classify_panel(section)
    return {
        "motivation": palette.accent,
        "method": palette.secondary,
        "benchmark": palette.warning,
        "innovation": palette.primary,
        "results": palette.success,
        "takeaway": palette.purple,
    }.get(kind, palette.primary)


def get_panel_figure(panel: Panel, task: PosterTask) -> Tuple[Optional[str], str]:
    if panel.figure:
        return panel.figure, panel.figure_caption or panel.figure_id
    if panel.figure_id and panel.figure_id in task.figures:
        fig = task.figures[panel.figure_id]
        return fig.image_source or fig.image_url, panel.figure_caption or fig.caption or panel.figure_id
    return None, panel.figure_caption or ""


def _normalise_label(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", (text or "").lower()).strip()


def _figure_target_keywords(fig: FigureAsset) -> List[str]:
    text = _normalise_label(" ".join([fig.best_matched_section, fig.type, fig.caption, fig.description]))
    targets: List[str] = []
    if any(k in text for k in ["method", "model", "architecture", "framework", "algorithm", "pipeline"]):
        targets.extend(["method", "framework", "approach"])
    if any(k in text for k in ["result", "experiment", "evaluation", "performance", "benchmark", "table"]):
        targets.extend(["result", "experiment", "evaluation", "benchmark"])
    if any(k in text for k in ["dataset", "data", "example", "case"]):
        targets.extend(["experiment", "benchmark", "result"])
    if any(k in text for k in ["finding", "analysis", "ablation"]):
        targets.extend(["finding", "analysis", "result"])
    return targets or _normalise_label(fig.best_matched_section).split()


def bind_available_figures(task: PosterTask, max_auto_figures: int = 2) -> None:
    """Attach unused source figures to matching CS poster panels.

    Dify/cached plans sometimes carry usable figure assets but leave every
    panel's ``figure_id`` empty. That makes the CS poster look text-only and
    also makes figure-reuse evaluation misleadingly zero. This deterministic
    binding is conservative: it only fills empty panel slots, never rewrites
    bullets, and caps the number of automatic figures.
    """

    if not task.figures or not task.panels:
        return
    used = {p.figure_id for p in task.panels if p.figure_id}
    candidates = []
    for fid, fig in task.figures.items():
        if fid in used:
            continue
        if (fig.audit_status or "").lower() != "ok":
            continue
        if not (fig.image_source or fig.image_url or fig.thumbnail_url):
            continue
        importance_rank = {"high": 0, "medium": 1, "low": 2}.get((fig.importance or "").lower(), 1)
        candidates.append((importance_rank, fid, fig))
    if not candidates:
        return
    candidates.sort(key=lambda item: (item[0], item[1]))

    assigned = 0
    for _, fid, fig in candidates:
        if assigned >= max_auto_figures:
            break
        targets = _figure_target_keywords(fig)
        direct_section = _normalise_label(fig.best_matched_section)
        best_panel: Optional[Panel] = None
        best_score = 0
        for panel in task.panels:
            if panel.figure_id or panel.figure:
                continue
            section = _normalise_label(panel.section)
            score = sum(1 for target in targets if target and target in section)
            if direct_section and (direct_section in section or section in direct_section):
                score += 3
            if score > best_score:
                best_panel = panel
                best_score = score
        if best_panel is None or best_score <= 0:
            continue
        best_panel.figure_id = fid
        best_panel.figure_caption = best_panel.figure_caption or fig.caption
        if best_panel.layout_hint == "text_only":
            best_panel.layout_hint = "text_left_image_right"
        assigned += 1


def fit_picture_box(img_w: int, img_h: int, x, y, w, h):
    img_ratio = img_w / max(img_h, 1)
    box_ratio = w / max(h, 1)
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = int(w / img_ratio)
    else:
        pic_h = h
        pic_w = int(h * img_ratio)
    return x + int((w - pic_w) / 2), y + int((h - pic_h) / 2), pic_w, pic_h


@lru_cache(maxsize=128)
def _image_aspect(image_source: str) -> Optional[float]:
    if not image_source:
        return None
    stream = load_image_from_source(image_source)
    if not stream:
        return None
    size = image_size(stream)
    if not size:
        return None
    return size[0] / max(size[1], 1)


def figure_squashed_in_vertical(image_source: str, box_w, box_h, min_height_ratio: float = 0.6) -> bool:
    """Whether a wide image would letterbox to < min_height_ratio of its
    allocated vertical-layout figure box. Used to decide whether to swap
    a text-top/image-bottom panel to a horizontal arrangement (or trim
    bullets to free vertical space) before committing the layout.

    Why: fit_picture_box centers a wider-than-box image by setting
    pic_h = box_w / img_ratio, so the picture occupies only
    box_ratio / img_ratio of the available height — a 2:1 image inside a
    1:1 figure box renders at 50% height even though the box looks full.
    """
    img_ratio = _image_aspect(image_source)
    if img_ratio is None:
        return False
    box_ratio = box_w / max(box_h, 1)
    if img_ratio <= box_ratio:
        return False
    return (box_ratio / img_ratio) < min_height_ratio



def _panel_font_size(base: float, panel: Optional[Panel] = None, task: Optional[PosterTask] = None) -> float:
    """Apply per-panel and global font scaling, clamped to a safe range.

    Mutated by the feedback loop: ``panel.body_font_scale`` shrinks when the
    panel is too crowded, grows when it has too much whitespace.
    ``task.global_font_scale`` does the same poster-wide. Clamps protect
    against runaway iterations producing unreadable text.
    """

    scale = 1.0
    if panel is not None and getattr(panel, "body_font_scale", None):
        scale *= max(0.7, min(1.3, panel.body_font_scale))
    if task is not None and getattr(task, "global_font_scale", None):
        scale *= max(0.8, min(1.2, task.global_font_scale))
    return base * scale


def add_figure(slide, x, y, w, h, image_source: str, caption: str, palette: Palette):
    add_rect(slide, x, y, w, h, RGBColor(249, 251, 253), palette.border, radius=True, line_width=0.6)
    try:
        img_stream = load_image_from_source(image_source)
        if not img_stream:
            add_textbox(slide, x, y + h / 2 - Inches(0.12), w, Inches(0.24), "No figure image", 8.5, palette.muted, False, PP_ALIGN.CENTER)
            return
        size = image_size(img_stream)
        px, py, pw, ph = x + Inches(0.05), y + Inches(0.05), w - Inches(0.10), h - Inches(0.10)
        if size:
            px, py, pw, ph = fit_picture_box(size[0], size[1], px, py, pw, ph)
        slide.shapes.add_picture(img_stream, px, py, width=pw, height=ph)
        if caption:
            add_textbox(slide, x, y + h - Inches(0.20), w, Inches(0.17), clean_text(caption, 80), 7.2, palette.muted, False, PP_ALIGN.CENTER)
    except Exception as exc:
        print(f"add_figure failed: {exc}")
        add_textbox(slide, x, y + h / 2 - Inches(0.12), w, Inches(0.24), "Figure load failed", 8.5, palette.muted, False, PP_ALIGN.CENTER)


def add_panel_header(slide, x, y, w, title: str, idx: int, palette: Palette, accent: RGBColor):
    bar_h = Inches(0.32)
    add_rect(slide, x, y, w, bar_h, palette.primary, radius=True)
    add_rect(slide, x, y, Inches(0.06), bar_h, accent, radius=False)

    add_textbox(slide, x + Inches(0.18), y + Inches(0.045), w - Inches(0.72), Inches(0.22), clean_text(title, 55), 12.5, palette.white, True, fit=True, min_font_size=8.8)

    badge = add_shape(slide, MSO_SHAPE.OVAL, x + w - Inches(0.34), y + Inches(0.055), Inches(0.20), Inches(0.20), accent)
    badge.line.fill.background()
    add_textbox(slide, x + w - Inches(0.34), y + Inches(0.055), Inches(0.20), Inches(0.20), str(idx), 7.5, palette.white, True, PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, h, panel: Panel, palette: Palette, accent: RGBColor, task: Optional[PosterTask] = None, max_items=5):
    items = panel.content[:max_items] or ["Key information will be filled by PlannerAgent."]
    row_h = h / max(len(items), 1)
    colors = [palette.primary, palette.success, palette.warning, palette.purple, palette.accent]
    body_size = _panel_font_size(8.8, panel, task)
    badge_size = _panel_font_size(7.2, panel, task)
    for idx, item in enumerate(items):
        item_y = y + row_h * idx
        c = colors[idx % len(colors)]
        size = Inches(0.20)
        dot = add_shape(slide, MSO_SHAPE.OVAL, x, item_y + Inches(0.045), size, size, c)
        dot.line.fill.background()
        add_textbox(slide, x, item_y + Inches(0.045), size, size, str(idx + 1), badge_size, palette.white, True, PP_ALIGN.CENTER, fit=True, min_font_size=5.6)
        text_h = min(max(Inches(0.42), row_h - Inches(0.02)), row_h)
        add_rich_textbox(
            slide,
            x + Inches(0.27),
            item_y + Inches(0.01),
            w - Inches(0.30),
            text_h,
            item,
            palette,
            font_size=body_size,
            bold=idx == 0,
            max_len=0,
            min_font_size=6.4,
        )


def add_goal_callout(slide, x, y, w, h, text: str, palette: Palette, task: Optional[PosterTask] = None, panel: Optional[Panel] = None):
    # Retained as a no-op for backward compatibility; the previous design
    # drew an alert-style "!" icon that read as an error indicator.
    return


def add_data_cards(slide, x, y, w, h, panel: Panel, palette: Palette, task: Optional[PosterTask] = None):
    numbers = []
    for item in panel.content:
        numbers.extend(re.findall(r"\d+(?:\.\d+)?%?|\d+\s?[xX×倍]", item))
    numbers = numbers[:3]
    if not numbers:
        return False

    gap = Inches(0.07)
    card_w = (w - gap * (len(numbers) - 1)) / len(numbers)
    metric_size = _panel_font_size(14, panel, task)
    label_size = _panel_font_size(7.3, panel, task)
    for idx, number in enumerate(numbers):
        cx = x + idx * (card_w + gap)
        add_rect(slide, cx, y, card_w, h, palette.soft, palette.border, radius=True, line_width=0.7)
        add_textbox(slide, cx + Inches(0.04), y + Inches(0.06), card_w - Inches(0.08), Inches(0.22), number, metric_size, palette.primary, True, PP_ALIGN.CENTER)
        add_textbox(slide, cx + Inches(0.04), y + Inches(0.31), card_w - Inches(0.08), Inches(0.18), "key metric", label_size, palette.muted, False, PP_ALIGN.CENTER)
    return True


def extract_metrics(panel: Panel, limit: int = 3) -> List[str]:
    metrics: List[str] = []
    for item in panel.content:
        for match in METRIC_RE.findall(item):
            if match not in metrics:
                metrics.append(match)
            if len(metrics) >= limit:
                return metrics
    return metrics


def add_metric_pills(slide, x, y, w, metrics: List[str], palette: Palette, accent: RGBColor):
    if not metrics:
        return
    gap = Inches(0.06)
    pill_w = min(Inches(0.58), (w - gap * (len(metrics) - 1)) / max(len(metrics), 1))
    for idx, metric in enumerate(metrics[:3]):
        px = x + idx * (pill_w + gap)
        add_rect(slide, px, y, pill_w, Inches(0.24), palette.soft, accent, radius=True, line_width=0.6)
        add_textbox(slide, px + Inches(0.02), y + Inches(0.03), pill_w - Inches(0.04), Inches(0.16), metric, 7.5, accent, True, PP_ALIGN.CENTER)


def add_section_glyph(slide, x, y, section: str, palette: Palette, accent: RGBColor):
    kind = classify_panel(section)
    glyph = {"motivation": "Q", "method": "M", "benchmark": "D", "results": "R", "takeaway": "T"}.get(kind, "S")
    add_shape(slide, MSO_SHAPE.OVAL, x, y, Inches(0.24), Inches(0.24), accent)
    add_textbox(slide, x, y + Inches(0.01), Inches(0.24), Inches(0.16), glyph, 8, palette.white, True, PP_ALIGN.CENTER)


def add_mini_pipeline(slide, x, y, w, h, panel: Panel, palette: Palette, accent: RGBColor, task: Optional[PosterTask] = None):
    steps = panel.content[:4] or ["Parse paper", "Plan panels", "Render poster", "Refine layout"]
    gap = Inches(0.10)
    box_w = (w - gap * (len(steps) - 1)) / len(steps)
    badge_size = _panel_font_size(8, panel, task)
    step_size = _panel_font_size(7.4, panel, task)
    for idx, step in enumerate(steps):
        sx = x + idx * (box_w + gap)
        add_rect(slide, sx, y, box_w, h, RGBColor(250, 252, 255), palette.border, radius=True, line_width=0.8)
        badge = add_shape(slide, MSO_SHAPE.OVAL, sx + box_w / 2 - Inches(0.11), y + Inches(0.08), Inches(0.22), Inches(0.22), accent)
        badge.line.fill.background()
        add_textbox(slide, sx + box_w / 2 - Inches(0.11), y + Inches(0.08), Inches(0.22), Inches(0.22), str(idx + 1), badge_size, palette.white, True, PP_ALIGN.CENTER, fit=True, min_font_size=5.8)
        add_rich_textbox(
            slide,
            sx + Inches(0.05),
            y + Inches(0.33),
            box_w - Inches(0.10),
            h - Inches(0.36),
            step,
            palette,
            font_size=step_size,
            bold=True,
            max_len=0,
            min_font_size=5.9,
        )
        if idx < len(steps) - 1:
            add_textbox(slide, sx + box_w - Inches(0.02), y + h / 2 - Inches(0.11), Inches(0.14), Inches(0.18), ">", 13, palette.text, True, PP_ALIGN.CENTER)


def add_panel_content(slide, x, y, w, h, panel: Panel, task: PosterTask, palette: Palette, idx: int):
    accent = panel_accent(panel.section, palette)
    add_rect(slide, x, y, w, h, palette.panel_bg, palette.border, radius=True, line_width=0.8)
    add_panel_header(slide, x, y, w, panel.section, idx, palette, accent)

    cx = x + Inches(0.13)
    cy = y + Inches(0.42)
    cw = w - Inches(0.26)
    ch = h - Inches(0.52)
    figure_source, figure_caption = get_panel_figure(panel, task)
    hint = panel.layout_hint or "text_only"
    kind = classify_panel(panel.section)

    if kind in ["benchmark", "results"] and not figure_source:
        if add_data_cards(slide, cx, cy, cw, Inches(0.58), panel, palette, task=task):
            add_bullets(slide, cx, cy + Inches(0.68), cw, ch - Inches(0.68), panel, palette, accent, task=task, max_items=4)
            return

    if figure_source and hint == "text_left_image_right":
        text_w = cw * 0.47
        add_bullets(slide, cx, cy, text_w, ch, panel, palette, accent, task=task, max_items=4)
        add_figure(slide, cx + text_w + Inches(0.12), cy + Inches(0.03), cw - text_w - Inches(0.12), ch - Inches(0.04), figure_source, figure_caption, palette)
    elif figure_source and hint == "image_focus":
        # Figure-dominant vertical: image takes ~75% of the panel and only the
        # single most-critical bullet sits below it. Applier sets this hint
        # when the VLM signals figure_too_small on a panel that's already
        # vertical, so the picture grows instead of staying letterboxed.
        fig_h = ch * 0.75
        text_h = ch - fig_h - Inches(0.10)
        add_figure(slide, cx, cy, cw, fig_h, figure_source, figure_caption, palette)
        add_bullets(slide, cx, cy + fig_h + Inches(0.10), cw, text_h, panel, palette, accent, task=task, max_items=1)
    elif figure_source and hint == "image_compact":
        # Compact figure container: figure takes 30% so the panel reads as
        # text-heavy with a small supporting visual. Applier uses this hint
        # to fix empty_space caused by oversized figure padding.
        fig_h = ch * 0.30
        text_h = ch - fig_h - Inches(0.10)
        add_bullets(slide, cx, cy, cw, text_h, panel, palette, accent, task=task, max_items=4)
        add_figure(slide, cx, cy + text_h + Inches(0.08), cw, fig_h, figure_source, figure_caption, palette)
    elif figure_source and hint in ["text_top_image_bottom", "image_top_text_bottom"]:
        # Probe whether vertical layout would crush the image into a thin
        # letterboxed strip. If so, prefer text_left_image_right when the
        # panel is wide enough; otherwise keep vertical but hand the figure
        # most of the height and reduce text to the single critical bullet.
        est_fig_h = ch * 0.56
        squashed = figure_squashed_in_vertical(figure_source, cw, est_fig_h)
        if squashed and cw >= Inches(3.4):
            text_w = cw * 0.47
            add_bullets(slide, cx, cy, text_w, ch, panel, palette, accent, task=task, max_items=3)
            add_figure(slide, cx + text_w + Inches(0.12), cy + Inches(0.03), cw - text_w - Inches(0.12), ch - Inches(0.04), figure_source, figure_caption, palette)
        elif squashed:
            fig_h = ch * 0.78
            if hint == "image_top_text_bottom":
                add_figure(slide, cx, cy, cw, fig_h, figure_source, figure_caption, palette)
                add_bullets(slide, cx, cy + fig_h + Inches(0.10), cw, ch - fig_h - Inches(0.10), panel, palette, accent, task=task, max_items=1)
            else:
                text_h = ch - fig_h - Inches(0.10)
                add_bullets(slide, cx, cy, cw, text_h, panel, palette, accent, task=task, max_items=1)
                add_figure(slide, cx, cy + text_h + Inches(0.08), cw, fig_h, figure_source, figure_caption, palette)
        elif hint == "image_top_text_bottom":
            fig_h = ch * 0.56
            add_figure(slide, cx, cy, cw, fig_h, figure_source, figure_caption, palette)
            add_bullets(slide, cx, cy + fig_h + Inches(0.15), cw, ch - fig_h - Inches(0.15), panel, palette, accent, task=task, max_items=3)
        else:
            text_h = ch * 0.40
            add_bullets(slide, cx, cy, cw, text_h, panel, palette, accent, task=task, max_items=3)
            add_figure(slide, cx, cy + text_h + Inches(0.08), cw, ch - text_h - Inches(0.10), figure_source, figure_caption, palette)
    elif figure_source and hint == "image_only":
        add_figure(slide, cx, cy, cw, ch, figure_source, figure_caption, palette)
    elif kind == "method":
        add_mini_pipeline(slide, cx, cy + Inches(0.02), cw, ch * 0.48, panel, palette, accent, task=task)
        add_bullets(slide, cx, cy + ch * 0.56, cw, ch * 0.42, panel, palette, accent, task=task, max_items=2)
    else:
        add_bullets(slide, cx, cy + Inches(0.02), cw, ch - Inches(0.02), panel, palette, accent, task=task, max_items=5)


class BasePosterTemplate:
    name = "base"

    def __init__(self, palette: Palette):
        self.palette = palette

    def render(self, task: PosterTask) -> io.BytesIO:
        raise NotImplementedError


class DashboardTemplate(BasePosterTemplate):
    name = "template_dashboard"

    def render(self, task: PosterTask) -> io.BytesIO:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        p = self.palette

        add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, p.bg, radius=False)
        self.add_header(slide, prs, task)

        body_x = Inches(0.10)
        body_y = Inches(0.86)
        body_w = prs.slide_width - Inches(0.20)
        body_h = prs.slide_height - Inches(1.34)
        grid = Grid(prs, body_x, body_y, body_w, body_h)
        spans = [
            grid.box(0, 0, 3, 4),
            grid.box(3, 0, 5, 4),
            grid.box(8, 0, 4, 4),
            grid.box(0, 4, 3, 4),
            grid.box(3, 4, 6, 4),
            grid.box(9, 4, 3, 4),
        ]

        for idx, panel in enumerate(sort_panels_for_dashboard(task.panels)[:6], start=1):
            pos = spans[idx - 1]
            add_panel_content(slide, pos["x"], pos["y"], pos["w"], pos["h"], panel, task, p, idx)

        self.add_footer(slide, prs, task)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    def add_header(self, slide, prs, task: PosterTask):
        p = self.palette
        add_rect(slide, 0, 0, prs.slide_width, Inches(0.82), p.navy, radius=False)
        title_w = prs.slide_width - Inches(3.38)
        title_text = clean_text(task.poster_title, 132)
        title_size = estimate_fit_font_size(title_text, title_w, Inches(0.48), max_size=21.0, min_size=13.5, bold=True)
        add_textbox(
            slide,
            Inches(0.34),
            Inches(0.06),
            title_w,
            Inches(0.48),
            title_text,
            title_size,
            p.white,
            True,
            fit=True,
            min_font_size=13.5,
        )
        subtitle = task.paper_info or "AI Agent workflow for paper-to-poster generation"
        add_textbox(slide, Inches(0.36), Inches(0.56), prs.slide_width - Inches(3.4), Inches(0.20), clean_text(subtitle, 120), 10.4, RGBColor(224, 236, 250), True, fit=True, min_font_size=7.5)
        add_textbox(slide, prs.slide_width - Inches(2.80), Inches(0.13), Inches(2.44), Inches(0.48), clean_text(task.authors or "Auto-generated Poster", 92), 10.3, p.white, False, PP_ALIGN.RIGHT, fit=True, min_font_size=7.0)

    def add_footer(self, slide, prs, task: PosterTask):
        p = self.palette
        y = prs.slide_height - Inches(0.44)
        add_rect(slide, Inches(0.10), y, prs.slide_width - Inches(0.20), Inches(0.34), p.primary, radius=True)
        add_textbox(slide, Inches(0.34), y + Inches(0.06), prs.slide_width - Inches(0.68), Inches(0.22), find_footer_sentence(task), 12, p.white, True, PP_ALIGN.CENTER)


class ClassicTemplate(DashboardTemplate):
    name = "template_classic"

    def render(self, task: PosterTask) -> io.BytesIO:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        p = self.palette

        add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, p.bg, radius=False)
        self.add_header(slide, prs, task)

        body_x = Inches(0.12)
        body_y = Inches(0.98)
        body_w = prs.slide_width - Inches(0.24)
        body_h = prs.slide_height - Inches(1.54)
        grid = Grid(prs, body_x, body_y, body_w, body_h, cols=12, rows=7, gap=Inches(0.10))
        spans = [
            grid.box(0, 0, 4, 4),
            grid.box(4, 0, 4, 4),
            grid.box(8, 0, 4, 4),
            grid.box(0, 4, 4, 3),
            grid.box(4, 4, 4, 3),
            grid.box(8, 4, 4, 3),
        ]

        add_textbox(slide, Inches(0.18), Inches(0.84), prs.slide_width - Inches(0.36), Inches(0.13), "Classic academic poster layout: balanced columns for scanning, comparison and discussion", 7.2, p.muted, False, PP_ALIGN.CENTER)

        for idx, panel in enumerate(sort_panels_for_dashboard(task.panels)[:6], start=1):
            pos = spans[idx - 1]
            add_panel_content(slide, pos["x"], pos["y"], pos["w"], pos["h"], panel, task, p, idx)

        self.add_footer(slide, prs, task)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

class StoryflowTemplate(DashboardTemplate):
    name = "template_storyflow"

    def render(self, task: PosterTask) -> io.BytesIO:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        p = self.palette

        add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, p.bg, radius=False)
        self.add_header(slide, prs, task)

        panels = sort_panels_for_dashboard(task.panels)[:7]
        self.add_story_ribbon(slide, prs)

        variant = self.resolve_variant(task, panels)
        if variant == "story_spotlight":
            self.render_spotlight(slide, prs, panels, task)
        elif variant == "story_zigzag":
            self.render_zigzag(slide, prs, panels, task)
        else:
            self.render_columns(slide, prs, panels[:6], task)

        self.add_footer(slide, prs, task)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    def resolve_variant(self, task: PosterTask, panels: List[Panel]) -> str:
        variant = getattr(task, "layout_variant", "auto") or "auto"
        if variant != "auto":
            return variant
        figure_count = sum(1 for p in panels if p.figure_id or p.figure)
        if figure_count >= 2:
            return "story_spotlight"
        if len(panels) != 6:
            return "story_zigzag"
        return "story_columns"

    def render_columns(self, slide, prs, panels: List[Panel], task: PosterTask):
        p = self.palette
        body_x = Inches(0.16)
        body_y = Inches(1.36)
        body_w = prs.slide_width - Inches(0.32)
        body_h = prs.slide_height - Inches(2.00)
        gap = Inches(0.09)
        card_w = (body_w - gap * (len(panels) - 1)) / max(len(panels), 1)

        for idx, panel in enumerate(panels, start=1):
            x = body_x + (idx - 1) * (card_w + gap)
            self.add_story_card(slide, x, body_y, card_w, body_h, panel, task, idx, compact=False)
            if idx < len(panels):
                add_textbox(slide, x + card_w - Inches(0.02), body_y + body_h / 2 - Inches(0.12), Inches(0.13), Inches(0.20), ">", 14, p.primary, True, PP_ALIGN.CENTER)

    def render_zigzag(self, slide, prs, panels: List[Panel], task: PosterTask):
        p = self.palette
        body_x = Inches(0.18)
        body_y = Inches(1.30)
        body_w = prs.slide_width - Inches(0.36)
        body_h = prs.slide_height - Inches(1.96)
        gap_x = Inches(0.12)
        gap_y = Inches(0.16)
        cols = min(4, max(3, (len(panels) + 1) // 2))
        card_w = (body_w - gap_x * (cols - 1)) / cols
        card_h = (body_h - gap_y) / 2
        for idx, panel in enumerate(panels[: cols * 2], start=1):
            r = 0 if idx <= cols else 1
            c = idx - 1 if r == 0 else cols * 2 - idx
            x = body_x + c * (card_w + gap_x)
            y = body_y + r * (card_h + gap_y)
            self.add_story_card(slide, x, y, card_w, card_h, panel, task, idx, compact=True)
            if idx < min(len(panels), cols * 2):
                add_textbox(slide, x + card_w - Inches(0.02), y + card_h / 2 - Inches(0.10), Inches(0.16), Inches(0.18), ">", 12, p.primary, True, PP_ALIGN.CENTER)

    def render_spotlight(self, slide, prs, panels: List[Panel], task: PosterTask):
        body_x = Inches(0.16)
        body_y = Inches(1.30)
        body_w = prs.slide_width - Inches(0.32)
        body_h = prs.slide_height - Inches(1.94)
        grid = Grid(prs, body_x, body_y, body_w, body_h, cols=12, rows=6, gap=Inches(0.11))
        spans = [
            grid.box(0, 0, 2, 3),
            grid.box(0, 3, 2, 3),
            grid.box(2, 0, 4, 3),
            grid.box(6, 0, 4, 3),
            grid.box(2, 3, 4, 3),
            grid.box(6, 3, 4, 3),
            grid.box(10, 0, 2, 6),
        ]
        for idx, panel in enumerate(panels[:7], start=1):
            pos = spans[idx - 1]
            self.add_story_card(slide, pos["x"], pos["y"], pos["w"], pos["h"], panel, task, idx, compact=pos["w"] < Inches(2.4))

    def add_story_ribbon(self, slide, prs):
        p = self.palette
        y = Inches(0.88)
        add_rect(slide, Inches(0.16), y, prs.slide_width - Inches(0.32), Inches(0.34), p.soft, p.border, radius=True, line_width=0.7)
        add_textbox(slide, Inches(0.32), y + Inches(0.06), Inches(2.2), Inches(0.20), "Storyflow", 12, p.primary, True)
        add_textbox(slide, Inches(2.00), y + Inches(0.065), prs.slide_width - Inches(2.42), Inches(0.18), "Problem framing -> method design -> evidence -> findings -> takeaways", 8.8, p.muted, False, PP_ALIGN.CENTER)

    def add_story_card(self, slide, x, y, w, h, panel: Panel, task: PosterTask, idx: int, compact: bool = False):
        p = self.palette
        accent = panel_accent(panel.section, p)
        add_rect(slide, x, y, w, h, p.panel_bg, p.border, radius=True, line_width=0.7)
        add_rect(slide, x, y, w, Inches(0.08), accent, radius=False)
        add_section_glyph(slide, x + Inches(0.10), y + Inches(0.16), panel.section, p, accent)

        marker = add_shape(slide, MSO_SHAPE.OVAL, x + w / 2 - Inches(0.18), y + Inches(0.18), Inches(0.36), Inches(0.36), accent)
        marker.line.fill.background()
        add_textbox(slide, x + w / 2 - Inches(0.18), y + Inches(0.18), Inches(0.36), Inches(0.36), str(idx), 14, p.white, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.08), y + Inches(0.62), w - Inches(0.16), Inches(0.44), clean_text(panel.section, 34), 10.3, p.primary, True, PP_ALIGN.CENTER)
        if getattr(task, "emphasis_level", "normal") in {"high", "strong"}:
            add_metric_pills(slide, x + Inches(0.12), y + Inches(0.98), w - Inches(0.24), extract_metrics(panel, 3), p, accent)

        figure_source, figure_caption = get_panel_figure(panel, task)
        content_y = y + (Inches(1.26) if getattr(task, "emphasis_level", "normal") in {"high", "strong"} else Inches(1.12))
        content_h = h - (content_y - y) - Inches(0.10)
        if figure_source:
            fig_ratio = 0.32 if compact or panel.layout_hint == "image_compact" else 0.42
            est_fig_w = w - Inches(0.20)
            est_fig_h = content_h * fig_ratio
            if figure_squashed_in_vertical(figure_source, est_fig_w, est_fig_h):
                fig_ratio = 0.70
            add_figure(slide, x + Inches(0.10), content_y, w - Inches(0.20), content_h * fig_ratio, figure_source, figure_caption, p)
            bullet_y = content_y + content_h * (fig_ratio + 0.06)
            bullets_squashed_cap = 1 if fig_ratio >= 0.6 else (3 if compact else 4)
            add_bullets(slide, x + Inches(0.10), bullet_y, w - Inches(0.20), content_h * (1 - fig_ratio - 0.08), panel, p, accent, task=task, max_items=bullets_squashed_cap)
        elif classify_panel(panel.section) == "method":
            add_mini_pipeline(slide, x + Inches(0.10), content_y + Inches(0.02), w - Inches(0.20), content_h * 0.40, panel, p, accent, task=task)
            add_bullets(slide, x + Inches(0.10), content_y + content_h * 0.48, w - Inches(0.20), content_h * 0.48, panel, p, accent, task=task, max_items=2 if compact else 3)
        else:
            add_bullets(slide, x + Inches(0.10), content_y, w - Inches(0.20), content_h, panel, p, accent, task=task, max_items=3 if compact else 4)

class MinimalTemplate(DashboardTemplate):
    name = "template_minimal"

    def render(self, task: PosterTask) -> io.BytesIO:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        p = self.palette

        add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(250, 251, 253), radius=False)
        self.add_minimal_header(slide, prs, task)

        body_x = Inches(0.42)
        body_y = Inches(1.20)
        body_w = prs.slide_width - Inches(0.84)
        body_h = prs.slide_height - Inches(1.82)
        grid = Grid(prs, body_x, body_y, body_w, body_h, cols=12, rows=6, gap=Inches(0.18))
        spans = [
            grid.box(0, 0, 4, 3),
            grid.box(4, 0, 4, 3),
            grid.box(8, 0, 4, 3),
            grid.box(0, 3, 4, 3),
            grid.box(4, 3, 4, 3),
            grid.box(8, 3, 4, 3),
        ]

        for idx, panel in enumerate(sort_panels_for_dashboard(task.panels)[:6], start=1):
            pos = spans[idx - 1]
            self.add_minimal_card(slide, pos["x"], pos["y"], pos["w"], pos["h"], panel, task, idx)

        y = prs.slide_height - Inches(0.42)
        add_textbox(slide, Inches(0.44), y, prs.slide_width - Inches(0.88), Inches(0.22), find_footer_sentence(task), 9.8, p.muted, True, PP_ALIGN.CENTER)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    def add_minimal_header(self, slide, prs, task: PosterTask):
        p = self.palette
        add_rect(slide, Inches(0.42), Inches(0.28), Inches(0.08), Inches(0.56), p.accent, radius=False)
        title_w = prs.slide_width - Inches(3.30)
        title_text = clean_text(task.poster_title, 132)
        title_size = estimate_fit_font_size(title_text, title_w, Inches(0.44), max_size=20.0, min_size=12.8, bold=True)
        add_textbox(slide, Inches(0.60), Inches(0.18), title_w, Inches(0.44), title_text, title_size, p.text, True, fit=True, min_font_size=12.8)
        add_textbox(slide, Inches(0.60), Inches(0.64), prs.slide_width - Inches(3.30), Inches(0.20), clean_text(task.paper_info or "Paper-to-poster summary", 120), 9.6, p.muted, False, fit=True, min_font_size=7.0)
        add_textbox(slide, prs.slide_width - Inches(2.80), Inches(0.26), Inches(2.30), Inches(0.42), clean_text(task.authors or "Auto-generated Poster", 92), 9.8, p.muted, False, PP_ALIGN.RIGHT, fit=True, min_font_size=6.8)
        add_rect(slide, Inches(0.42), Inches(0.96), prs.slide_width - Inches(0.84), Inches(0.02), p.border, radius=False)

    def add_minimal_card(self, slide, x, y, w, h, panel: Panel, task: PosterTask, idx: int):
        p = self.palette
        accent = panel_accent(panel.section, p)
        add_rect(slide, x, y, w, h, p.white, RGBColor(220, 226, 234), radius=True, line_width=0.6)
        add_rect(slide, x + Inches(0.16), y + Inches(0.15), Inches(0.34), Inches(0.05), accent, radius=False)
        add_textbox(slide, x + Inches(0.16), y + Inches(0.24), w - Inches(0.54), Inches(0.30), clean_text(panel.section, 45), 12.2, p.text, True)
        add_textbox(slide, x + w - Inches(0.42), y + Inches(0.21), Inches(0.22), Inches(0.22), f"{idx:02d}", 8.2, p.muted, True, PP_ALIGN.RIGHT)

        figure_source, figure_caption = get_panel_figure(panel, task)
        cx = x + Inches(0.16)
        cy = y + Inches(0.65)
        cw = w - Inches(0.32)
        ch = h - Inches(0.82)
        if figure_source:
            fig_ratio = 0.78 if figure_squashed_in_vertical(figure_source, cw, ch * 0.48) else 0.48
            add_figure(slide, cx, cy, cw, ch * fig_ratio, figure_source, figure_caption, p)
            bullet_top = cy + ch * (fig_ratio + 0.06)
            bullet_h = ch * (1 - fig_ratio - 0.06)
            add_bullets(slide, cx, bullet_top, cw, bullet_h, panel, p, accent, max_items=1 if fig_ratio >= 0.6 else 3)
        else:
            add_bullets(slide, cx, cy, cw, ch, panel, p, accent, max_items=4)


def find_footer_sentence(task: PosterTask) -> str:
    for panel in task.panels:
        if classify_panel(panel.section) == "takeaway" and panel.content:
            return clean_text(panel.content[0], 170)
    if task.panels and task.panels[-1].content:
        return clean_text(task.panels[-1].content[0], 170)
    return "This poster summarizes the motivation, method, results and key takeaways of the paper."


def choose_palette(task: PosterTask) -> Palette:
    return PALETTES.get(task.color_theme or "", PALETTES["academic_blue"])


def choose_template(task: PosterTask, palette: Palette) -> BasePosterTemplate:
    template_map = {
        "template_dashboard": DashboardTemplate,
        "template_classic": ClassicTemplate,
        "template_storyflow": StoryflowTemplate,
        "template_minimal": MinimalTemplate,
    }
    cls = template_map.get(task.template or "template_dashboard", DashboardTemplate)
    return cls(palette)


def generate_dashboard_pptx(task: PosterTask) -> io.BytesIO:
    bind_available_figures(task)
    palette = choose_palette(task)
    template = choose_template(task, palette)
    return template.render(task)
