import copy
import json
import shutil
import subprocess
import tempfile
import time
import uuid
from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont

from app.models import Panel, PosterTask
from app.ppt_renderer import generate_dashboard_pptx
from app.vlm_commenter import check_layout_with_vlm


ISSUE_OVERLAPPING_ELEMENTS = "overlapping_elements"
ISSUE_EMPTY_SPACE = "empty_space"
ISSUE_LOW_CONTRAST = "low_contrast"
ISSUE_FIGURE_TOO_SMALL = "figure_too_small"

# Closed action set understood by FeedbackApplier. Mirrors
# ``SVFPSuggestedAction`` in vlm_commenter.py so the schema and the dispatcher
# stay in lockstep.
ACTION_REDUCE_BULLET_COUNT = "reduce_bullet_count"
ACTION_SHRINK_TEXT = "shrink_text"
ACTION_TRUNCATE_BULLETS = "truncate_bullets"
ACTION_SHRINK_FIGURE_BOX = "shrink_figure_box"
ACTION_ENLARGE_FONT = "enlarge_font"
ACTION_ADD_BULLET = "add_bullet"
ACTION_COMPACT_FIGURE_BOX = "compact_figure_box"
ACTION_SWITCH_PALETTE = "switch_palette"
ACTION_NONE = "none"

# Palette rotation for ``switch_palette``. Cycles through the two
# high-contrast themes so each iteration produces a visually distinct
# poster — escapes the academic_blue → low_contrast → academic_blue
# deadlock from the previous design.
PALETTE_CYCLE = ["academic_blue", "engineering_green"]
HIGH_CONTRAST_THEMES = set(PALETTE_CYCLE)
STORYFLOW_VARIANTS = ["story_columns", "story_spotlight", "story_zigzag"]


@dataclass
class PanelFeedback:
    section: str
    issues: List[str] = field(default_factory=list)
    suggested_action: str = "none"
    target_value: Optional[Any] = None


@dataclass
class LayoutFeedback:
    score: float = 7.0
    global_issues: List[str] = field(default_factory=list)
    panel_feedback: List[PanelFeedback] = field(default_factory=list)
    comment: str = ""
    source: str = "heuristic"

    def to_dict(self) -> Dict[str, Any]:
        return {
            "score": self.score,
            "global_issues": self.global_issues,
            "panel_feedback": [pf.__dict__ for pf in self.panel_feedback],
            "comment": self.comment,
            "source": self.source,
        }


class PreviewRenderer:
    """Fast structural preview for feedback before generating the final PPTX.

    Designed so that small content changes between iterations are visible
    in the rendered PNG:

    * bullets wrap to multiple lines instead of being clipped at 60 chars;
    * the panel header shows the *real* bullet count;
    * a "iter N · panels=M · bullets=Σ" stamp is drawn so two iterations
      with different content never produce a byte-identical image;
    * bullets that overflow the panel box are drawn with a red marker so
      the VLM can see overflow visually rather than guessing.
    """

    def __init__(self, width: int = 1600, height: int = 900):
        self.width = width
        self.height = height
        self.font_title = self._font(36, bold=True)
        self.font_panel = self._font(18, bold=True)
        self.font_body = self._font(14)
        self.font_small = self._font(12)

    def _font(self, size: int, bold: bool = False):
        candidates = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
            "/Library/Fonts/Arial Unicode.ttf",
        ]
        for path in candidates:
            if Path(path).exists():
                try:
                    return ImageFont.truetype(path, size=size)
                except Exception:
                    pass
        return ImageFont.load_default()

    def _text_width(self, draw: ImageDraw.ImageDraw, text: str, font) -> int:
        try:
            bbox = draw.textbbox((0, 0), text, font=font)
            return bbox[2] - bbox[0]
        except Exception:
            return len(text) * 7

    def _wrap_text(self, draw: ImageDraw.ImageDraw, text: str, max_width: int, font) -> List[str]:
        """Greedy character-level wrap that respects CJK + Latin mix."""

        text = " ".join((text or "").split())
        if not text:
            return [""]
        lines: List[str] = []
        buf = ""
        for ch in text:
            candidate = buf + ch
            if self._text_width(draw, candidate, font) <= max_width:
                buf = candidate
                continue
            if buf:
                lines.append(buf)
                buf = ch
            else:
                lines.append(ch)
                buf = ""
        if buf:
            lines.append(buf)
        return lines

    def render(self, task: PosterTask, output_path: Path, iteration: Optional[int] = None) -> Path:
        img = Image.new("RGB", (self.width, self.height), "#f4f8fd")
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, self.width, 92), fill="#002b75")
        draw.text((42, 22), task.poster_title[:70], fill="white", font=self.font_title)
        draw.text((42, 62), (task.paper_info or "Poster preview")[:110], fill="#dcecff", font=self.font_body)
        draw.text((self.width - 300, 34), (task.authors or "Auto-generated")[:42], fill="white", font=self.font_body)

        panels = task.panels[:6]
        boxes = self._boxes(len(panels))
        total_bullets = 0
        for idx, panel in enumerate(panels):
            x, y, w, h = boxes[idx]
            draw.rounded_rectangle((x, y, x + w, y + h), radius=10, fill="white", outline="#b1c7e2", width=2)
            draw.rectangle((x, y, x + w, y + 38), fill="#0052ae")
            header = f"{idx + 1}. {panel.section[:30]}  ({len(panel.content)} bullets)"
            draw.text((x + 14, y + 9), header, fill="white", font=self.font_panel)

            cy = y + 50
            body_bottom = y + h - 18
            text_max_width = w - 60
            for bullet_idx, bullet in enumerate(panel.content):
                total_bullets += 1
                lines = self._wrap_text(draw, bullet, text_max_width, self.font_body) or [""]
                bullet_height = max(28, 4 + 18 * len(lines))
                overflowed = cy + bullet_height > body_bottom
                marker_fill = "#dc2626" if overflowed else "#2286ef"
                draw.ellipse((x + 18, cy + 2, x + 38, cy + 22), fill=marker_fill)
                draw.text((x + 24, cy + 2), str(bullet_idx + 1), fill="white", font=self.font_body)
                ly = cy
                for line in lines:
                    draw.text((x + 48, ly), line, fill="#142036", font=self.font_body)
                    ly += 18
                cy += bullet_height
                if overflowed:
                    draw.line((x + 8, body_bottom + 4, x + w - 8, body_bottom + 4), fill="#dc2626", width=2)
                    draw.text(
                        (x + 14, body_bottom + 6),
                        f"overflow: {len(panel.content) - bullet_idx - 1} more hidden",
                        fill="#dc2626",
                        font=self.font_small,
                    )
                    break

            if panel.figure_id:
                draw.rounded_rectangle((x + w - 150, y + h - 82, x + w - 18, y + h - 20), radius=7, fill="#e8f1fc", outline="#b1c7e2")
                draw.text((x + w - 136, y + h - 58), panel.figure_id, fill="#0052ae", font=self.font_panel)

        stamp = f"iter={iteration if iteration is not None else '-'} · panels={len(task.panels)} · bullets={total_bullets} · theme={task.color_theme}"
        draw.text((42, self.height - 26), stamp, fill="#345", font=self.font_small)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(output_path)
        return output_path

    def _boxes(self, count: int):
        margin = 20
        top = 110
        gap = 14
        body_w = self.width - margin * 2
        body_h = self.height - top - 74
        if count <= 4:
            cols, rows = 2, 2
        else:
            cols, rows = 3, 2
        w = (body_w - gap * (cols - 1)) // cols
        h = (body_h - gap * (rows - 1)) // rows
        boxes = []
        for r in range(rows):
            for c in range(cols):
                boxes.append((margin + c * (w + gap), top + r * (h + gap), w, h))
        return boxes


class HeuristicLayoutChecker:
    """Rule-based fallback when VLM is unavailable or returns garbage.

    Only emits the 3 SVFP issues (overlapping_elements / empty_space /
    low_contrast). The heuristic is intentionally conservative — it produces
    smaller scores and fewer issues than a real VLM call, but enough to keep
    the closed loop making progress when offline.
    """

    def check(self, task: PosterTask, source: str = "heuristic") -> LayoutFeedback:
        feedback = LayoutFeedback(score=8.0, source=source)
        total_penalty = 0.0

        # Per-panel: a panel is "overlapping" if it has too many bullets or
        # any bullet far exceeds the rendered text-box width; it's "empty"
        # if it has very little content.
        for panel in task.panels:
            bullet_lengths = [len(item) for item in panel.content]
            n_bullets = len(panel.content)
            total_chars = sum(bullet_lengths)

            if n_bullets > 5 or any(length > 120 for length in bullet_lengths):
                feedback.panel_feedback.append(
                    PanelFeedback(
                        section=panel.section,
                        issues=[ISSUE_OVERLAPPING_ELEMENTS],
                        suggested_action=ACTION_REDUCE_BULLET_COUNT,
                        target_value=4,
                    )
                )
                total_penalty += 0.6
            elif n_bullets <= 1 or total_chars < 30:
                feedback.panel_feedback.append(
                    PanelFeedback(
                        section=panel.section,
                        issues=[ISSUE_EMPTY_SPACE],
                        suggested_action=ACTION_ENLARGE_FONT,
                        target_value=1.15,
                    )
                )
                total_penalty += 0.4

        # Global signal: many panels each holding >5 bullets means the
        # whole poster reads as crowded.
        crowded_panels = sum(1 for p in task.panels if len(p.content) > 5)
        if crowded_panels >= 2:
            feedback.global_issues.append(ISSUE_OVERLAPPING_ELEMENTS)
            total_penalty += 0.5

        feedback.score = max(1.0, min(10.0, round(8.5 - total_penalty, 1)))
        feedback.comment = "Heuristic layout check completed."
        return feedback


class FeedbackApplier:
    """Translate :class:`LayoutFeedback` into concrete edits on a ``PosterTask``.

    Each ``suggested_action`` maps to exactly one mutation primitive. The
    set is closed (see ``ACTION_*`` constants); unknown actions fall through
    to a small issue-based default so a malformed VLM reply still makes
    progress. ``low_contrast`` always cycles to the next palette in
    :data:`PALETTE_CYCLE` — this breaks the old no-op trap where
    ``academic_blue`` was already considered "high contrast" so nothing
    changed even when the VLM kept reporting the issue.
    """

    # Conservative bounds so runaway iterations can't produce unreadable
    # posters. These mirror the clamps in ppt_renderer._panel_font_size.
    FONT_SCALE_MIN = 0.7
    FONT_SCALE_MAX = 1.3
    GLOBAL_FONT_SCALE_MIN = 0.8
    GLOBAL_FONT_SCALE_MAX = 1.2

    def apply(self, task: PosterTask, feedback: LayoutFeedback) -> PosterTask:
        next_task = copy.deepcopy(task)
        by_section = {panel.section: panel for panel in next_task.panels}
        touched: set[str] = set()

        for item in feedback.panel_feedback:
            panel = by_section.get(item.section)
            if not panel:
                continue
            applied = self._dispatch_panel_action(panel, item, next_task)
            if applied:
                touched.add(panel.section)

        # Global issues handled after the per-panel pass so panel-level
        # fixes are visible when we decide global remediation.
        gi = set(feedback.global_issues)
        if ISSUE_LOW_CONTRAST in gi:
            self._rotate_palette(next_task)
        if ISSUE_EMPTY_SPACE in gi:
            self._enlarge_global_font(next_task)
            self._increase_visual_emphasis(next_task)
            self._rotate_layout_variant(next_task)
        if ISSUE_OVERLAPPING_ELEMENTS in gi:
            # If multiple panels overlap, shrink the global font as well —
            # the per-panel pass already trimmed/reduced the worst offenders.
            self._shrink_global_font(next_task)

        return next_task

    # ------------------------------------------------------------------
    # Dispatch
    # ------------------------------------------------------------------

    def _dispatch_panel_action(self, panel: Panel, item: "PanelFeedback", task: PosterTask) -> bool:
        action = (item.suggested_action or "").strip()
        target = item.target_value
        issues = list(item.issues or [])
        horizontal = panel.layout_hint == "text_left_image_right"
        has_figure = bool(panel.figure_id or panel.figure)

        # figure_too_small only makes sense for vertical layouts where the
        # picture is actually letterboxed. On a healthy horizontal panel the
        # image isn't squashed — reducing bullets just leaves whitespace.
        # Strip the issue (and any action it dragged in) before dispatch.
        if ISSUE_FIGURE_TOO_SMALL in issues and horizontal:
            issues = [i for i in issues if i != ISSUE_FIGURE_TOO_SMALL]
            if action == ACTION_REDUCE_BULLET_COUNT and (target == 1 or target == "1"):
                action = ACTION_NONE
            if not issues and action == ACTION_NONE:
                return False

        # figure_too_small on a vertical layout: drop to one critical bullet
        # AND flip the hint to image_focus so the figure grows instead of
        # staying compact. Direct route bypasses the noisy action mapping.
        if ISSUE_FIGURE_TOO_SMALL in issues and has_figure and not horizontal:
            self._reduce_bullets(panel, 1)
            panel.layout_hint = "image_focus"
            return True

        if action == ACTION_REDUCE_BULLET_COUNT:
            count = int(target) if target else 4
            # Guard against the empty-panel regression: don't strip a
            # horizontal panel down to a single bullet unless we have a real
            # overlap signal forcing it.
            if horizontal and count <= 1 and ISSUE_OVERLAPPING_ELEMENTS not in issues:
                return False
            self._reduce_bullets(panel, count)
            return True
        if action == ACTION_SHRINK_TEXT:
            self._scale_panel_font(panel, float(target) if target else 0.85)
            return True
        if action == ACTION_TRUNCATE_BULLETS:
            self._truncate_bullets(panel, int(target) if target else 80)
            return True
        if action == ACTION_SHRINK_FIGURE_BOX or action == ACTION_COMPACT_FIGURE_BOX:
            # Don't demote a working horizontal panel into image_compact —
            # image_compact is text-dominant (figure shrinks to 30%) and was
            # the root cause of the iter-1→iter-2 regression where Method
            # and Experiments lost their wide image columns. Also keep
            # hands off image_focus, which was set deliberately to grow the
            # figure; switching back would unwind that.
            if horizontal:
                self._scale_panel_font(panel, 0.92)
            elif panel.layout_hint == "image_focus":
                self._scale_panel_font(panel, 1.05)
            else:
                panel.layout_hint = "image_compact"
            return True
        if action == ACTION_ENLARGE_FONT:
            self._scale_panel_font(panel, float(target) if target else 1.15)
            return True
        if action == ACTION_ADD_BULLET:
            # Don't fabricate filler content unrelated to the paper —
            # adding meaningless bullets to balance density is worse than
            # leaving the panel sparse. Treat empty_space by enlarging the
            # existing text instead.
            self._scale_panel_font(panel, 1.15)
            return True
        if action == ACTION_SWITCH_PALETTE:
            self._rotate_palette(task)
            return True
        if action == ACTION_NONE or not action:
            # No explicit action — pick a sensible default from the issues.
            return self._apply_issue_defaults(panel, issues, task)

        # Unknown action string — treat it like "none" rather than crashing.
        return self._apply_issue_defaults(panel, issues, task)

    # ------------------------------------------------------------------
    # Mutation primitives
    # ------------------------------------------------------------------

    def _reduce_bullets(self, panel: Panel, target_count: int) -> None:
        panel.content = [self._shorten(item) for item in panel.content[: max(1, target_count)]]

    def _truncate_bullets(self, panel: Panel, max_len: int) -> None:
        panel.content = [self._shorten(b, max_len) for b in panel.content]

    def _scale_panel_font(self, panel: Panel, factor: float) -> None:
        new_scale = max(
            self.FONT_SCALE_MIN,
            min(self.FONT_SCALE_MAX, panel.body_font_scale * factor),
        )
        panel.body_font_scale = round(new_scale, 3)

    def _enlarge_global_font(self, task: PosterTask) -> None:
        task.global_font_scale = round(
            min(self.GLOBAL_FONT_SCALE_MAX, task.global_font_scale * 1.10), 3
        )

    def _increase_visual_emphasis(self, task: PosterTask) -> None:
        task.emphasis_level = "high"

    def _shrink_global_font(self, task: PosterTask) -> None:
        task.global_font_scale = round(
            max(self.GLOBAL_FONT_SCALE_MIN, task.global_font_scale * 0.92), 3
        )

    def _rotate_palette(self, task: PosterTask) -> None:
        """Cycle through PALETTE_CYCLE so VLM sees a visually different image
        next iteration. Breaks the academic_blue ↔ low_contrast deadlock by
        always moving — never no-op."""

        current = task.color_theme or PALETTE_CYCLE[0]
        try:
            idx = PALETTE_CYCLE.index(current)
            next_idx = (idx + 1) % len(PALETTE_CYCLE)
        except ValueError:
            next_idx = 0
        task.color_theme = PALETTE_CYCLE[next_idx]

    def _rotate_layout_variant(self, task: PosterTask) -> None:
        if task.template != "template_storyflow":
            return
        current = task.layout_variant or "auto"
        if current == "auto":
            task.layout_variant = "story_spotlight"
            return
        try:
            idx = STORYFLOW_VARIANTS.index(current)
            task.layout_variant = STORYFLOW_VARIANTS[(idx + 1) % len(STORYFLOW_VARIANTS)]
        except ValueError:
            task.layout_variant = "story_spotlight"

    def _apply_issue_defaults(self, panel: Panel, issues: List[str], task: PosterTask) -> bool:
        """Best-effort fix when the VLM provides issues but no action."""

        horizontal = panel.layout_hint == "text_left_image_right"
        has_figure = bool(panel.figure_id or panel.figure)

        if ISSUE_FIGURE_TOO_SMALL in issues:
            # On horizontal panels the image isn't actually squashed, so
            # reducing bullets would just empty the text column. Skip and
            # fall through to remaining issues (if any).
            if not horizontal and has_figure:
                self._reduce_bullets(panel, 1)
                panel.layout_hint = "image_focus"
                return True
        if ISSUE_OVERLAPPING_ELEMENTS in issues:
            self._reduce_bullets(panel, 4)
            return True
        if ISSUE_EMPTY_SPACE in issues:
            self._scale_panel_font(panel, 1.15)
            return True
        if ISSUE_LOW_CONTRAST in issues:
            self._rotate_palette(task)
            return True
        return False

    def _shorten(self, text: str, max_len: int = 92) -> str:
        text = " ".join((text or "").split())
        if len(text) <= max_len:
            return text
        return text[: max_len - 3] + "..."


_SOFFICE_FALLBACK_PATHS = (
    # macOS Homebrew --cask install location (not in PATH by default).
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    # Common Linux locations when the binary isn't symlinked.
    "/usr/bin/soffice",
    "/usr/lib/libreoffice/program/soffice",
)

# 2400x1350 PNG instead of soffice's ~1024-wide default. The VLM can't read
# panel headers and bullets at default resolution, so the screenshot branch
# effectively gives up before the model even gets a chance to comment.
_PNG_EXPORT_FILTER = (
    'png:impress_png_Export:'
    '{"PixelWidth":{"type":"long","value":"2400"},'
    '"PixelHeight":{"type":"long","value":"1350"}}'
)

# Cached preflight outcome for the current process. ``None`` = not probed
# yet; ``False`` = soffice is broken and ``render_pptx_to_png`` fast-fails
# instead of paying the 90s timeout on every iteration.
_PREFLIGHT_RESULT: Optional[bool] = None


def _find_soffice() -> Optional[str]:
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    for path in _SOFFICE_FALLBACK_PATHS:
        if Path(path).exists():
            return path
    return None


def _kill_stale_soffice() -> None:
    """Force-kill any leftover soffice processes. Best-effort — pkill returns
    1 when nothing matches, which we ignore."""
    try:
        subprocess.run(
            ["pkill", "-9", "-f", "soffice"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=5,
        )
    except Exception:
        pass


def _run_soffice_once(
    soffice: str, pptx_path: Path, output_dir: Path
) -> Tuple[int, str]:
    """One conversion attempt with an isolated profile."""
    try:
        with tempfile.TemporaryDirectory(prefix="lo_profile_") as profile_dir:
            profile_uri = Path(profile_dir).resolve().as_uri()
            result = subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    "--nodefault",
                    "--norestore",
                    f"-env:UserInstallation={profile_uri}",
                    "--convert-to",
                    _PNG_EXPORT_FILTER,
                    "--outdir",
                    str(output_dir),
                    str(pptx_path),
                ],
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=90,
            )
            stderr = (result.stderr or b"").decode(errors="ignore")[:400]
            return result.returncode, stderr
    except subprocess.TimeoutExpired:
        return -9, f"timeout after 90s on {pptx_path.name}"
    except Exception as exc:
        return -1, f"unexpected error: {exc}"


def render_pptx_to_png(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    if _PREFLIGHT_RESULT is False:
        return None

    soffice = _find_soffice()
    if not soffice:
        print(
            "[render_pptx_to_png] LibreOffice not found. Install with "
            "`brew install --cask libreoffice` (macOS) or "
            "`apt install libreoffice` (Linux)."
        )
        return None

    output_dir.mkdir(parents=True, exist_ok=True)
    rc, stderr = _run_soffice_once(soffice, pptx_path, output_dir)

    # -6 = SIGABRT. Usually a stale background soffice holding a profile lock
    # or a first-launch Gatekeeper abort. Killing leftover processes plus a
    # fresh temp profile (already isolated per attempt) clears both.
    if rc == -6:
        print(
            f"[render_pptx_to_png] soffice SIGABRT on {pptx_path.name}; "
            "killing stale soffice and retrying once"
        )
        _kill_stale_soffice()
        time.sleep(0.5)
        rc, stderr = _run_soffice_once(soffice, pptx_path, output_dir)

    if rc != 0:
        print(f"[render_pptx_to_png] soffice exit={rc}: {stderr}")
        return None

    candidates = sorted(output_dir.glob(f"{pptx_path.stem}*.png"))
    if not candidates:
        print(f"[render_pptx_to_png] soffice ran but no PNG appeared in {output_dir}")
        return None
    return candidates[0]


def preflight_soffice() -> bool:
    """Probe soffice once per process with a minimal PPTX, then cache.

    Without this, a permanently broken soffice burns 90s × iterations
    waiting for timeouts. With it, we set ``_PREFLIGHT_RESULT = False`` and
    every subsequent ``render_pptx_to_png`` call returns ``None`` instantly.
    """
    global _PREFLIGHT_RESULT
    if _PREFLIGHT_RESULT is not None:
        return _PREFLIGHT_RESULT

    if not _find_soffice():
        _PREFLIGHT_RESULT = False
        return False

    try:
        from pptx import Presentation

        with tempfile.TemporaryDirectory(prefix="lo_preflight_") as tmp:
            tmp_path = Path(tmp)
            probe = tmp_path / "probe.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "soffice preflight"
            prs.save(str(probe))
            png = render_pptx_to_png(probe, tmp_path)
            _PREFLIGHT_RESULT = png is not None and png.exists()
    except Exception as exc:
        print(f"[preflight_soffice] probe failed: {exc}")
        _PREFLIGHT_RESULT = False

    if not _PREFLIGHT_RESULT:
        print(
            "[preflight_soffice] soffice cannot render PNG; VLM screenshot "
            "feedback disabled for this session. On macOS try: "
            "xattr -dr com.apple.quarantine /Applications/LibreOffice.app"
        )
    return _PREFLIGHT_RESULT


def parse_vlm_feedback(raw: Dict[str, Any], fallback: LayoutFeedback) -> LayoutFeedback:
    """Convert a raw VLM payload into a :class:`LayoutFeedback`.

    Policy:

    * ``disabled`` / ``vlm_error`` → return the heuristic ``fallback`` so
      the loop keeps making progress.
    * ``vlm`` / ``vlm_partial`` → trust the VLM. ``vlm_partial`` is the
      truncation-recovered case: we still extract whichever fields are
      present and prefer them over the heuristic, because the VLM saw the
      actual rendered preview.
    * ``vlm_unparsed`` → final resort: try one more time to dig a JSON out
      of the comment field; only fall back to heuristics if that also fails.
    """

    if not raw:
        return fallback

    source = str(raw.get("source", ""))
    if source in {"disabled", "vlm_error"}:
        fallback.comment = (fallback.comment + " " + str(raw.get("comment", ""))).strip()
        return fallback

    data: Optional[Dict[str, Any]] = None
    if "score" in raw or "panel_feedback" in raw or "global_issues" in raw:
        data = raw
    else:
        comment = str(raw.get("comment", ""))
        try:
            start = comment.index("{")
            end = comment.rindex("}") + 1
            data = json.loads(comment[start:end])
        except Exception:
            data = None

    if not data:
        fallback.comment = (
            fallback.comment + " VLM output unparseable: " + str(raw.get("comment", ""))[:200]
        ).strip()
        fallback.source = "vlm_unparsed_fallback"
        return fallback

    panel_feedback: List[PanelFeedback] = []
    for item in data.get("panel_feedback", []) or []:
        if not isinstance(item, dict):
            continue
        panel_feedback.append(
            PanelFeedback(
                section=str(item.get("section", "")),
                issues=[str(i) for i in (item.get("issues") or [])],
                suggested_action=str(item.get("suggested_action", "none")),
                target_value=item.get("target_value"),
            )
        )

    try:
        score = float(data.get("score", fallback.score))
    except (TypeError, ValueError):
        score = fallback.score
    score = max(1.0, min(10.0, score))

    return LayoutFeedback(
        score=score,
        global_issues=[str(g) for g in (data.get("global_issues") or [])],
        panel_feedback=panel_feedback,
        comment=str(data.get("comment", fallback.comment)).strip()
        or fallback.comment
        or "VLM returned no comment.",
        source=str(data.get("source", source or "vlm")),
    )


class VisualFeedbackLoop:
    def __init__(self):
        self.preview_renderer = PreviewRenderer()
        self.heuristic_checker = HeuristicLayoutChecker()
        self.applier = FeedbackApplier()

    def run(self, task: PosterTask) -> Dict[str, Any]:
        from app.run_archive import RunArchive, update_runs_index

        # Probe soffice once before the loop starts. If it's broken (missing
        # binary, Gatekeeper quarantine, persistent SIGABRT), the cached
        # negative result makes every render_pptx_to_png call below return
        # None instantly instead of waiting on a 90s timeout per iteration.
        preflight_soffice()

        run_id = uuid.uuid4().hex[:12]
        archive = RunArchive.create(run_id, task.poster_title)
        archive.save_input(task.model_dump())

        current = copy.deepcopy(task)
        max_iterations = max(1, min(current.max_iterations or 1, 20))
        best_task = copy.deepcopy(current)
        best_score = -1.0
        best_issue_count: Optional[int] = None
        history: List[Dict[str, Any]] = []
        stop_reason = "max_iterations_reached"
        converged = False

        for iteration in range(1, max_iterations + 1):
            iter_started_at = datetime.now(timezone.utc).isoformat(timespec="seconds")

            # Draft preview is kept as a debug artifact so users can eyeball
            # structural changes between iterations, but it is NOT fed to
            # the VLM — that role belongs exclusively to the real PPT
            # screenshot below.
            preview_path = self.preview_renderer.render(
                current,
                archive.preview_dir / f"iter_{iteration}_preview.png",
                iteration=iteration,
            )

            pptx_buf = generate_dashboard_pptx(current)
            pptx_path = archive.pptx_dir / f"iter_{iteration}.pptx"
            pptx_path.write_bytes(pptx_buf.getvalue())
            screenshot_path = render_pptx_to_png(pptx_path, archive.pptx_dir)

            if screenshot_path:
                screenshot_fallback = self.heuristic_checker.check(
                    current, source="ppt_screenshot_heuristic"
                )
                try:
                    image = Image.open(screenshot_path).convert("RGB")
                    vlm_raw = check_layout_with_vlm(image)
                    feedback = parse_vlm_feedback(vlm_raw, screenshot_fallback)
                except Exception as exc:
                    feedback = screenshot_fallback
                    feedback.source = "vlm_error_fallback"
                    feedback.comment = f"{feedback.comment} VLM call failed: {exc}"
            else:
                # LibreOffice unavailable — degrade to heuristics rather
                # than feeding a fake draft to the VLM. The source label
                # makes this obvious in logs and the run report.
                feedback = self.heuristic_checker.check(
                    current, source="heuristic_no_screenshot"
                )
                feedback.comment = (
                    f"{feedback.comment} LibreOffice unavailable; "
                    "install with `brew install --cask libreoffice` to enable real VLM feedback."
                )

            iter_finished_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
            record = {
                "iteration": iteration,
                "started_at": iter_started_at,
                "finished_at": iter_finished_at,
                "score": feedback.score,
                "feedback": feedback.to_dict(),
                "task_snapshot": current.model_dump(),
                "preview_image": str(preview_path),
                "pptx": str(pptx_path),
                "ppt_screenshot": str(screenshot_path) if screenshot_path else "",
            }
            history.append(record)

            current_issue_count = len(feedback.global_issues) + sum(
                len(pf.issues) for pf in feedback.panel_feedback
            )
            score_improved = feedback.score > best_score

            tie_with_design_change = (
                feedback.score == best_score
                and best_issue_count is not None
                and current_issue_count <= best_issue_count
                and (
                    getattr(current, "layout_variant", "auto") != getattr(best_task, "layout_variant", "auto")
                    or getattr(current, "color_theme", "") != getattr(best_task, "color_theme", "")
                    or getattr(current, "emphasis_level", "normal") != getattr(best_task, "emphasis_level", "normal")
                    or getattr(current, "global_font_scale", 1.0) != getattr(best_task, "global_font_scale", 1.0)
                )
            )

            if score_improved or tie_with_design_change:
                best_score = feedback.score
                best_task = copy.deepcopy(current)
                best_issue_count = current_issue_count
            elif best_issue_count is None or current_issue_count < best_issue_count:
                best_issue_count = current_issue_count

            if not feedback.global_issues and not feedback.panel_feedback:
                converged = True
                stop_reason = "no_issues"
                break

            current = self.applier.apply(current, feedback)
        else:
            stop_reason = "max_iterations_reached"

        final_buf = generate_dashboard_pptx(best_task)
        final_path = archive.save_final_pptx_bytes(final_buf.getvalue())

        summary = {
            "best_score": best_score,
            "iterations": len(history),
            "max_iterations": max_iterations,
            "converged": converged,
            "convergence_reason": stop_reason,
            "score_curve": [r["score"] for r in history],
            "issue_curve": [
                len(r["feedback"].get("global_issues") or [])
                + sum(len(pf.get("issues") or []) for pf in (r["feedback"].get("panel_feedback") or []))
                for r in history
            ],
            "vlm_sources": [r["feedback"].get("source") for r in history],
            "best_task_snapshot": best_task.model_dump(),
        }
        run_report_path = archive.save_report(
            input_task=task.model_dump(),
            summary=summary,
            iterations=history,
        )
        try:
            update_runs_index()
        except Exception as exc:  # never let index regeneration take down a run
            print(f"update_runs_index failed: {exc}")

        return {
            "run_id": run_id,
            "best_score": best_score,
            "iterations": len(history),
            "converged": converged,
            "convergence_reason": stop_reason,
            "history": history,
            "final_path": str(final_path),
            "final_filename": final_path.name,
            "task": best_task,
            "archive_dir": str(archive.run_dir),
            "run_folder": archive.folder_name,
            "run_report_path": str(run_report_path),
        }


# ---------------------------------------------------------------------------
# Convergence detection (closed-loop iteration controller)
# ---------------------------------------------------------------------------


@dataclass
class ConvergenceConfig:
    """Configuration knobs for :class:`ConvergenceDetector`.

    Attributes
    ----------
    max_iterations:
        Hard upper bound on the number of iterations. Reaching it forces
        ``converged=True`` with reason ``max_iterations_reached``.
    excellent_threshold:
        Score above which the loop is considered "good enough" and stops
        immediately.
    min_delta:
        Minimum absolute score improvement between consecutive iterations
        that still counts as progress. Anything below this is "stagnant".
    stagnant_patience:
        Number of consecutive stagnant rounds tolerated before stopping.
    adaptive:
        If ``True``, ``min_delta`` is auto-tightened as the score climbs so
        that high scores require larger improvements to keep iterating.
    min_iterations:
        Always run at least this many iterations regardless of other
        conditions. Useful for ablation studies.
    """

    max_iterations: int = 5
    excellent_threshold: float = 9.0
    min_delta: float = 0.2
    stagnant_patience: int = 2
    adaptive: bool = True
    min_iterations: int = 1


@dataclass
class ConvergenceState:
    """Snapshot returned by :meth:`ConvergenceDetector.update`."""

    converged: bool
    reason: str
    iteration: int
    score: float
    best_score: float
    delta: float
    stagnant_rounds: int
    effective_min_delta: float

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


class ConvergenceDetector:
    """Closed-loop convergence detector for the visual feedback loop.

    The detector consumes one iteration result at a time via
    :meth:`update` and returns a :class:`ConvergenceState` flagging whether
    the loop should stop. It supports four stop conditions:

    1. ``excellent_threshold`` — score is high enough to ship.
    2. ``max_iterations`` — iteration budget exhausted.
    3. ``stagnant_patience`` — score plateaued for too many rounds.
    4. ``no_issues`` — when the caller reports an empty issue list.

    Example
    -------
    >>> detector = ConvergenceDetector(ConvergenceConfig(max_iterations=3))
    >>> detector.update(score=7.2, feedback={"global_issues": ["dense_content"]}).converged
    False
    >>> detector.update(score=9.4, feedback={"global_issues": []}).converged
    True
    """

    def __init__(self, config: Optional[ConvergenceConfig] = None) -> None:
        self.config = config or ConvergenceConfig()
        self._iteration: int = 0
        self._best_score: float = float("-inf")
        self._last_score: Optional[float] = None
        self._stagnant_rounds: int = 0
        self._history: List[ConvergenceState] = []

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _adaptive_min_delta(self, current_score: float) -> float:
        """Adaptively tighten ``min_delta`` as the score climbs.

        At score ≥ 9 we expect very small improvements; at low scores we
        accept larger ones because there is more headroom. Returns the
        original ``min_delta`` when ``adaptive`` is disabled.
        """

        base = self.config.min_delta
        if not self.config.adaptive:
            return base
        if current_score >= 9.0:
            return max(base * 0.5, 0.05)
        if current_score >= 8.0:
            return max(base * 0.75, 0.1)
        return base

    @staticmethod
    def _issue_count(feedback: Optional[Dict[str, Any]]) -> int:
        """Count total issues reported in a feedback dict.

        Tolerates both the legacy ``LayoutFeedback`` shape (``global_issues``
        + ``panel_feedback``) and the new SVFP list shape.
        """

        if not feedback:
            return 0
        if isinstance(feedback, list):
            return len(feedback)
        total = 0
        total += len(feedback.get("global_issues", []) or [])
        for panel in feedback.get("panel_feedback", []) or []:
            total += len(panel.get("issues", []) or [])
        return total

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def update(
        self,
        score: float,
        feedback: Optional[Dict[str, Any]] = None,
    ) -> ConvergenceState:
        """Feed one iteration's result and obtain a convergence verdict.

        Parameters
        ----------
        score:
            VLM (or heuristic) score for this iteration, expected in [1, 10].
        feedback:
            Full feedback payload for this iteration. Used to detect the
            ``no_issues`` short-circuit. May be ``None``.

        Returns
        -------
        ConvergenceState
            The updated state, including a ``converged`` flag and a
            human-readable ``reason``.
        """

        self._iteration += 1
        previous = self._last_score
        delta = 0.0 if previous is None else score - previous
        self._last_score = score

        effective_delta = self._adaptive_min_delta(score)
        improved = previous is None or delta >= effective_delta
        if improved:
            self._stagnant_rounds = 0
        else:
            self._stagnant_rounds += 1

        if score > self._best_score:
            self._best_score = score

        converged, reason = self._decide(score, feedback, effective_delta)

        state = ConvergenceState(
            converged=converged,
            reason=reason,
            iteration=self._iteration,
            score=score,
            best_score=self._best_score,
            delta=delta,
            stagnant_rounds=self._stagnant_rounds,
            effective_min_delta=effective_delta,
        )
        self._history.append(state)
        return state

    def _decide(
        self,
        score: float,
        feedback: Optional[Dict[str, Any]],
        effective_delta: float,
    ) -> Tuple[bool, str]:
        """Apply stop conditions in priority order and return (flag, reason)."""

        cfg = self.config
        if self._iteration < cfg.min_iterations:
            return False, "below_min_iterations"
        if score >= cfg.excellent_threshold:
            return True, "excellent_threshold"
        if self._issue_count(feedback) == 0 and self._iteration > 0:
            return True, "no_issues"
        if self._stagnant_rounds >= cfg.stagnant_patience:
            return True, f"stagnant_patience(delta<{effective_delta:.3f})"
        if self._iteration >= cfg.max_iterations:
            return True, "max_iterations_reached"
        return False, "in_progress"

    # ------------------------------------------------------------------
    # Inspection helpers
    # ------------------------------------------------------------------

    @property
    def history(self) -> List[ConvergenceState]:
        """Return a copy of all accumulated states."""

        return list(self._history)

    @property
    def best_score(self) -> float:
        """Return the best score observed so far (``-inf`` until first update)."""

        return self._best_score

    @property
    def iteration(self) -> int:
        """Number of :meth:`update` calls made on this detector."""

        return self._iteration

    def reset(self) -> None:
        """Clear all internal counters so the detector can be reused."""

        self._iteration = 0
        self._best_score = float("-inf")
        self._last_score = None
        self._stagnant_rounds = 0
        self._history.clear()


# ---------------------------------------------------------------------------
# Batch-experiment friendly interface
# ---------------------------------------------------------------------------


def check_convergence(
    iterations: List[Dict[str, Any]],
    config: Optional[ConvergenceConfig] = None,
) -> Dict[str, Any]:
    """Replay a list of iteration records through a fresh detector.

    Designed for batch experiment scripts that collect raw iteration logs
    and need a single function to decide whether the loop converged.

    Parameters
    ----------
    iterations:
        Ordered list of dicts. Each dict must contain a ``score`` field and
        optionally a ``feedback`` field. Records produced by
        :class:`VisualFeedbackLoop` (which use the key ``feedback``) are
        accepted directly.
    config:
        Optional :class:`ConvergenceConfig`. Defaults are used when omitted.

    Returns
    -------
    dict
        ``{"converged": bool, "reason": str, "best_score": float,
        "iterations_used": int, "states": [ConvergenceState.to_dict, ...]}``
    """

    detector = ConvergenceDetector(config)
    final_state: Optional[ConvergenceState] = None
    for record in iterations:
        score = float(record.get("score", 0.0))
        feedback = record.get("feedback") or record.get("feedback_list")
        state = detector.update(score=score, feedback=feedback)
        final_state = state
        if state.converged:
            break

    if final_state is None:
        return {
            "converged": False,
            "reason": "empty_input",
            "best_score": float("-inf"),
            "iterations_used": 0,
            "states": [],
        }

    return {
        "converged": final_state.converged,
        "reason": final_state.reason,
        "best_score": detector.best_score,
        "iterations_used": detector.iteration,
        "states": [s.to_dict() for s in detector.history],
    }


def _convergence_demo() -> None:
    """Quick demo printing convergence states for a synthetic run."""

    fake_run = [
        {"score": 6.4, "feedback": {"global_issues": ["dense_content"], "panel_feedback": []}},
        {"score": 7.6, "feedback": {"global_issues": ["dense_content"], "panel_feedback": []}},
        {"score": 7.7, "feedback": {"global_issues": [], "panel_feedback": []}},
    ]
    print(json.dumps(check_convergence(fake_run), ensure_ascii=False, indent=2))


if __name__ == "__main__":
    _convergence_demo()
