"""Shared pptx geometry walker used by B1 (layout) and B2 (readability).

Extracts:
* panel boxes — large rectangles that bound each panel card
* text runs — font size + character count per text frame
* slide dimensions

Implementations are intentionally simple-and-defensible: deterministic,
no LLM call, no random sampling. Edge cases (single-panel poster,
empty text frames) fall through to neutral values rather than raising
— B1/B2 already short-circuit on those.
"""

from __future__ import annotations

import statistics
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


__all__ = [
    "PanelBox",
    "TextSpan",
    "extract_panel_boxes",
    "extract_text_spans",
    "slide_dimensions_emu",
]


@dataclass(frozen=True)
class PanelBox:
    """A panel card on the poster, in EMU (English Metric Units)."""
    x: int
    y: int
    w: int
    h: int

    @property
    def area(self) -> int:
        return self.w * self.h

    @property
    def x2(self) -> int:
        return self.x + self.w

    @property
    def y2(self) -> int:
        return self.y + self.h


@dataclass(frozen=True)
class TextSpan:
    """One paragraph's worth of text and its font size (pt)."""
    text: str
    font_pt: float
    container_area_emu: int   # area of the text frame's container; 0 if unknown


# Heuristics for telling a "panel card" rectangle apart from background /
# banner / decoration. The poster templates put panel cards as AUTO_SHAPE
# rectangles with area roughly 1.5–25 in² (smaller = decoration; larger =
# full-slide background or banner). These thresholds were calibrated on
# the dashboard / classic / minimal templates.
_MIN_PANEL_AREA_IN2 = 1.5
_BACKGROUND_WIDTH_FRACTION = 0.90    # spans ≥90% of slide width AND
_BACKGROUND_HEIGHT_FRACTION = 0.50   # ≥50% of slide height → background
_BANNER_HEIGHT_FRACTION = 0.20       # spans nearly full width but short → banner


def slide_dimensions_emu(pptx_path: Path) -> Tuple[int, int]:
    prs = Presentation(str(pptx_path))
    return int(prs.slide_width or 0), int(prs.slide_height or 0)


def extract_panel_boxes(pptx_path: Path) -> List[PanelBox]:
    """Return panel-card rectangles from the first slide.

    Filters out:
    * the slide-wide background rectangle
    * the top banner that holds the poster title
    * tiny decorative AUTO_SHAPEs (bullet markers, dividers)
    """
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return []
    slide = prs.slides[0]
    slide_w = int(prs.slide_width or 0)
    slide_h = int(prs.slide_height or 0)

    boxes: List[PanelBox] = []
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        x = int(sh.left or 0)
        y = int(sh.top or 0)
        w = int(sh.width or 0)
        h = int(sh.height or 0)
        if w <= 0 or h <= 0:
            continue
        area_in2 = Emu(w).inches * Emu(h).inches
        if area_in2 < _MIN_PANEL_AREA_IN2:
            continue
        # Full-slide background
        if slide_w > 0 and w >= slide_w * _BACKGROUND_WIDTH_FRACTION \
                and h >= slide_h * _BACKGROUND_HEIGHT_FRACTION:
            continue
        # Top banner (wide but short)
        if slide_w > 0 and w >= slide_w * _BACKGROUND_WIDTH_FRACTION \
                and h < slide_h * _BANNER_HEIGHT_FRACTION:
            continue
        boxes.append(PanelBox(x=x, y=y, w=w, h=h))
    return boxes


def extract_text_spans(pptx_path: Path) -> List[TextSpan]:
    """Return one ``TextSpan`` per non-empty paragraph in the first slide.

    Font size defaults to 18pt when the run inherits from layout/master
    (python-pptx returns ``None`` in that case). This matches the
    visually-perceived default the dashboard renderer uses.
    """
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return []
    slide = prs.slides[0]

    spans: List[TextSpan] = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        container_w = int(sh.width or 0)
        container_h = int(sh.height or 0)
        container_area = container_w * container_h
        for para in sh.text_frame.paragraphs:
            text = "".join(run.text or "" for run in para.runs)
            if not text.strip():
                continue
            font_pt: Optional[float] = None
            for run in para.runs:
                size = run.font.size
                if size is not None:
                    font_pt = size.pt
                    break
            if font_pt is None and para.font.size is not None:
                font_pt = para.font.size.pt
            if font_pt is None:
                font_pt = 18.0  # template default
            spans.append(TextSpan(text=text, font_pt=float(font_pt), container_area_emu=container_area))
    return spans


# ---------------------------------------------------------------------------
# Small numeric helpers shared by B1 / B2
# ---------------------------------------------------------------------------


def mad_normalised(values: List[float]) -> float:
    """Median absolute deviation, normalised by the median (CV-like).
    Returns 0 when fewer than 2 samples or median is zero."""
    if len(values) < 2:
        return 0.0
    med = statistics.median(values)
    if med == 0:
        return 0.0
    mad = statistics.median(abs(v - med) for v in values)
    return mad / abs(med)
