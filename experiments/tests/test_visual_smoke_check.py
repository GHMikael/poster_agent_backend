from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from experiments.metrics.base import MetricContext, MetricRegistry
import experiments.metrics.visual_smoke_check  # noqa: F401


def _ctx(tmp_path: Path, pptx_path: Path) -> MetricContext:
    cell = tmp_path / "ours_svfp_paper"
    cell.mkdir(exist_ok=True)
    return MetricContext(
        artifact_dir=cell,
        pptx_path=pptx_path,
        png_path=None,
        panels_json=None,
        experiment_log_path=None,
        paper_path=Path("/no/such.pdf"),
        paper_meta={},
        config={},
    )


def _save_pptx(path: Path, *, text: str, width=4.0, height=0.4, font=18):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font)
    prs.save(path)


def test_visual_smoke_penalizes_likely_text_overflow(tmp_path):
    pptx = tmp_path / "overflow.pptx"
    _save_pptx(
        pptx,
        text="This is a very long technical title that cannot fit inside a tiny one-line text box",
        width=2.0,
        height=0.18,
        font=22,
    )

    result = MetricRegistry.get("visual_smoke_check")().compute(_ctx(tmp_path, pptx))

    assert result.score < 1.0
    assert result.extra["n_likely_overflow"] >= 1


def test_visual_smoke_accepts_reasonable_text_box(tmp_path):
    pptx = tmp_path / "ok.pptx"
    _save_pptx(pptx, text="Readable method summary", width=4.0, height=0.6, font=14)

    result = MetricRegistry.get("visual_smoke_check")().compute(_ctx(tmp_path, pptx))

    assert result.score == 1.0
    assert result.extra["n_likely_overflow"] == 0
