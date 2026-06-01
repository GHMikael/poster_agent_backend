"""Cheap structural smoke check for rendered PPTX quality.

This metric catches issues that are painful to discover after a full n=30 run:
text boxes that are likely too small, excessive ellipses, tiny fonts, and top
banner text overlaps. It is a preflight guard, not a substitute for B1/B2 or
human visual evaluation.
"""

from __future__ import annotations

import math
import re
from dataclasses import dataclass
from typing import List, Tuple

from pptx import Presentation

from experiments.metrics.base import Metric, MetricContext, MetricResult, MetricRegistry


EMU_PER_INCH = 914400.0


@dataclass(frozen=True)
class _TextBox:
    text: str
    font_pt: float
    x: int
    y: int
    w: int
    h: int

    @property
    def x2(self) -> int:
        return self.x + self.w

    @property
    def y2(self) -> int:
        return self.y + self.h

    @property
    def area(self) -> int:
        return self.w * self.h


def _font_size(shape, para) -> float:
    for run in para.runs:
        if run.font.size is not None:
            return float(run.font.size.pt)
    if para.font.size is not None:
        return float(para.font.size.pt)
    return 12.0


def _extract_text_boxes(pptx_path) -> List[_TextBox]:
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return []
    boxes: List[_TextBox] = []
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip()).strip()
        if not text:
            continue
        font_pt = max(_font_size(shape, p) for p in shape.text_frame.paragraphs if p.text.strip())
        boxes.append(
            _TextBox(
                text=re.sub(r"\s+", " ", text),
                font_pt=font_pt,
                x=int(shape.left or 0),
                y=int(shape.top or 0),
                w=int(shape.width or 0),
                h=int(shape.height or 0),
            )
        )
    return boxes


def _likely_text_overflow(box: _TextBox, *, slack: float = 1.18) -> bool:
    text = box.text.strip()
    if len(text) <= 6 or box.w <= 0 or box.h <= 0:
        return False
    width_pt = box.w / EMU_PER_INCH * 72.0
    height_pt = box.h / EMU_PER_INCH * 72.0
    avg_char_width = max(box.font_pt * 0.52, 1.0)
    chars_per_line = max(int(width_pt / avg_char_width), 1)
    longest_word = max((len(t) for t in re.split(r"\s+", text)), default=1)
    lines = max(math.ceil(len(text) / chars_per_line), math.ceil(longest_word / chars_per_line))
    required = lines * box.font_pt * 1.15
    return required > height_pt * slack


def _overlap_ratio(a: _TextBox, b: _TextBox) -> float:
    x0, y0 = max(a.x, b.x), max(a.y, b.y)
    x1, y1 = min(a.x2, b.x2), min(a.y2, b.y2)
    if x1 <= x0 or y1 <= y0:
        return 0.0
    inter = (x1 - x0) * (y1 - y0)
    return inter / max(min(a.area, b.area), 1)


@MetricRegistry.register
class VisualSmokeCheck(Metric):
    metric_id = "visual_smoke_check"
    description = "Preflight guard for text overflow, ellipsis pressure and top-banner overlap."

    def compute(self, ctx: MetricContext) -> MetricResult:
        if not ctx.config.get("enabled", True):
            return self._skip("disabled in metrics.yaml")
        if not ctx.pptx_path.exists():
            return self._skip(f"pptx missing: {ctx.pptx_path}")

        boxes = _extract_text_boxes(ctx.pptx_path)
        if not boxes:
            return self._skip("no text boxes")

        overflow = [b.text[:90] for b in boxes if _likely_text_overflow(b)]
        ellipsis = [b.text[:90] for b in boxes if "..." in b.text or "…" in b.text]
        tiny_fonts = [b.text[:90] for b in boxes if len(b.text) > 8 and b.font_pt < 6.0]
        top_boxes = [b for b in boxes if b.y < 0.95 * EMU_PER_INCH and len(b.text) > 3]
        top_overlaps: List[Tuple[str, str]] = []
        for i, a in enumerate(top_boxes):
            for b in top_boxes[i + 1 :]:
                if _overlap_ratio(a, b) > 0.12:
                    top_overlaps.append((a.text[:60], b.text[:60]))

        penalty = (
            0.14 * len(overflow)
            + 0.035 * len(ellipsis)
            + 0.12 * len(tiny_fonts)
            + 0.18 * len(top_overlaps)
        )
        score = max(0.0, 1.0 - min(1.0, penalty))
        return MetricResult(
            metric_id=self.metric_id,
            score=score,
            extra={
                "n_text_boxes": len(boxes),
                "n_likely_overflow": len(overflow),
                "n_ellipsis": len(ellipsis),
                "n_tiny_fonts": len(tiny_fonts),
                "n_top_overlaps": len(top_overlaps),
                "overflow_examples": overflow[:6],
                "ellipsis_examples": ellipsis[:6],
                "top_overlap_examples": top_overlaps[:3],
            },
            notes="Preflight smoke metric; use rendered PNG inspection and B1/B2 for final visual claims.",
        )
